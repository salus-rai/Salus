"""
# SPDX-License-Identifier: MIT
# Copyright 2024 - 2025 Infosys Ltd.

Permission is hereby granted, free of charge, to any person obtaining a copy of this software and associated documentation files (the "Software"), to deal in the Software without restriction, including without limitation the rights to use, copy, modify, merge, publish, distribute, sublicense, and/or sell copies of the Software, and to permit persons to whom the Software is furnished to do so, subject to the following conditions:
 
The above copyright notice and this permission notice shall be included in all copies or substantial portions of the Software.
 
THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY, FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM, OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE SOFTWARE.
"""

import base64
import io
from typing import Tuple
from PIL import Image
from privacy.service.imagePrivacy import AttributeDict, ImagePrivacy
from privacy.service.api_req import ApiCall
from privacy.service.textPrivacy import TextPrivacy
from privacy.service.__init__ import *
import numpy as np
from privacy.config.logger import request_id_var
from privacy.config.logger import CustomLogger
import os
import threading
import uuid
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from unidecode import unidecode
import tempfile

log = CustomLogger()
error_dict = {}

class PPTService:
    def processImages(slide, shape, payload, uid):
        """
        Process the images in a PowerPoint slide.

        Args:
            slide (PowerPoint slide): The slide object in which the image is located.
            shape (PowerPoint shape): The shape object representing the image.
            payload (dict): The payload containing the necessary information for image processing.
            uid (str): The unique identifier for the image.

        Returns:
            None

        Raises:
            Exception: If an error occurs during image processing.

        """
        try:
            request_id_var.set(uid)
            image = shape.image
            img_bytes = image.blob
            if len(img_bytes) < 700:
                return None
            imgd = io.BytesIO(img_bytes)
            payload["image"] = AttributeDict({"file": imgd})
            # payload["piiEntitiesToBeRedacted"] = None
            resImage = ImagePrivacy.image_anonymize(AttributeDict(payload))
            resImg = base64.b64decode(resImage)
            img_stream = io.BytesIO(resImg)
            img = Image.open(img_stream)
            
            # Use tempfile to create a temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_file:
                img_path = tmp_file.name
                img.save(img_path)
            
            slide.shapes.add_picture(img_path, shape.left, shape.top, shape.width, shape.height)
            slide.shapes._spTree.remove(shape._element)
            os.remove(img_path)
        except Exception as e:
            log.error(str(e))
            log.error("Line No:" + str(e.__traceback__.tb_lineno))
            log.error(str(e.__traceback__.tb_frame))
            if request_id_var.get() not in error_dict:
                error_dict[request_id_var.get()] = []
            error_dict[request_id_var.get()].append({"UUID": request_id_var.get(), "function": "PPTMASKMainFunction", "msg": str(e.__class__.__name__), "description": str(e) + "Line No:" + str(e.__traceback__.tb_lineno)})
            raise Exception(e)

    def editText(text, i, shape):
        """
        Replaces a specific portion of the text in a PowerPoint shape with a placeholder tag.

        Args:
            text (str): The original text in the shape.
            i (slice): The slice object representing the portion of the text to be replaced.
            shape (PowerPoint.Shape): The PowerPoint shape object.

        Returns:
            None
        """
        request_id_var.set("editText")
        log.debug(str(text[i.start:i.end]) + ":" + str(i.entity_type))
        shape.text = shape.text.replace(text[i.start:i.end], f"<{i.entity_type}>")

    def processText(slide, payload, uid):
        """
        Process the text in a slide by analyzing and anonymizing it.

        Args:
            slide (Slide): The slide object containing the text to be processed.
            payload (Payload): The payload object containing the necessary data for processing.
            uid (str): The unique identifier for the request.

        Raises:
            ValueError: If the result from textAnalyze is not a list.
            Exception: If any other error occurs during the processing.

        Returns:
            None
        """
        try:
            request_id_var.set(uid)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = unidecode(shape.text)
                     
                    accDetails = None
                    if payload.portfolio is not None:
                        accDetails = AttributeDict({"portfolio": payload.portfolio, "account": payload.account})
                    res = TextPrivacy.textAnalyze(text=text, accName=accDetails, exclusion=payload.exclusion.split(',') if payload.exclusion is not None else [],piiEntitiesToBeRedacted=payload.piiEntitiesToBeRedacted.split(",") if payload.piiEntitiesToBeRedacted is not None else [],nlp=payload.nlp,scoreThreshold=payload.scoreThreshold if payload.scoreThreshold != None else 0.0)
                     
            
                    if not isinstance(res, list):
                        raise ValueError(f"Expected a list from textAnalyze, but got: {type(res)}")
                    res = anonymizer._remove_conflicts_and_get_text_manipulation_data(res, (ConflictResolutionStrategy.MERGE_SIMILAR_OR_CONTAINED))
                    res = anonymizer._merge_entities_with_whitespace_between(text, res)
                    resThreads = []
                    for i in res:
                        thread = threading.Thread(target=PPTService.editText, args=(text, i, shape))
                        thread.start()
                        resThreads.append(thread)
                    for thread in resThreads:
                        thread.join()
        except Exception as e:
            log.error(str(e))
            log.error("Line No:" + str(e.__traceback__.tb_lineno))
            log.error(str(e.__traceback__.tb_frame))
            if request_id_var.get() not in error_dict:
                error_dict[request_id_var.get()] = []
            error_dict[request_id_var.get()].append({"UUID": request_id_var.get(), "function": "PPTMASKMainFunction", "msg": str(e.__class__.__name__), "description": str(e) + "Line No:" + str(e.__traceback__.tb_lineno)})
            raise Exception(e)

    def processTables(slide, payload, uid):
        """
        Process the tables in a PowerPoint slide and apply text privacy analysis and redaction.

        Args:
            slide (PowerPoint slide): The slide to process.
            payload (object): The payload containing the necessary data for text privacy analysis.
            uid (str): The unique identifier for the request.

        Raises:
            Exception: If an error occurs during the processing.

        Returns:
            None
        """
        try:
            request_id_var.set(uid)
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            text = unidecode(cell.text)+" ,"
                            accDetails = None
                            if payload.portfolio is not None:
                                accDetails = AttributeDict({"portfolio": payload.portfolio, "account": payload.account})
                            res = TextPrivacy.textAnalyze(text=text, accName=accDetails, exclusion=payload.exclusion.split(',') if payload.exclusion is not None else [],piiEntitiesToBeRedacted=payload.piiEntitiesToBeRedacted.split(",") if payload.piiEntitiesToBeRedacted is not None else [],nlp=payload.nlp,scoreThreshold=payload.scoreThreshold if payload.scoreThreshold != None else 0.0)
                            res = anonymizer._remove_conflicts_and_get_text_manipulation_data(res, (ConflictResolutionStrategy.MERGE_SIMILAR_OR_CONTAINED))
                            res = anonymizer._merge_entities_with_whitespace_between(text, res)
                            for i in res:
                                cell.text = cell.text.replace(text[i.start:i.end], f"<{i.entity_type}>")
        except Exception as e:
            log.error(str(e))
            log.error("Line No:" + str(e.__traceback__.tb_lineno))
            log.error(str(e.__traceback__.tb_frame))
            if request_id_var.get() not in error_dict:
                error_dict[request_id_var.get()] = []
            error_dict[request_id_var.get()].append({"UUID": request_id_var.get(), "function": "PPTMASKMainFunction", "msg": str(e.__class__.__name__), "description": str(e) + "Line No:" + str(e.__traceback__.tb_lineno)})
            raise Exception(e)

    def mask_ppt(payload):
        """
        Masks the content of a PowerPoint presentation.

        Args:
            payload (dict): The payload containing the PowerPoint presentation and other information.

        Returns:
            io.BytesIO: The masked PowerPoint presentation as a BytesIO object.

        Raises:
            Exception: If an error occurs during the masking process.
        """
        try:
            log.debug("payload:-" + str(payload))
            id = uuid.uuid4().hex
            request_id_var.set(id)
            payload = AttributeDict(payload)
            if payload.portfolio is not None or payload.account is not None:
                response_value = ApiCall.request(AttributeDict({"portfolio": payload.portfolio, "account": payload.account}))
                if response_value is None:
                    return None

            ppt_file = Presentation(io.BytesIO(payload.file.file.read()))

            for slide in ppt_file.slides:
                threads = []
                thread = threading.Thread(target=PPTService.processText, args=(slide, payload, id))
                thread.start()
                threads.append(thread)
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        thread = threading.Thread(target=PPTService.processImages, args=(slide, shape, payload, id))
                        thread.start()
                        threads.append(thread)
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        thread = threading.Thread(target=PPTService.processTables, args=(slide, payload, id))
                        thread.start()
                        threads.append(thread)
                for thread in threads:
                    thread.join()

            ppt_bytes = io.BytesIO()
            ppt_file.save(ppt_bytes)
            ppt_bytes.seek(0)

            return ppt_bytes
        except Exception as e:
            log.error(str(e))
            log.error("Line No:" + str(e.__traceback__.tb_lineno))
            log.error(str(e.__traceback__.tb_frame))
            if request_id_var.get() not in error_dict:
                error_dict[request_id_var.get()] = []
            error_dict[request_id_var.get()].append({"UUID": request_id_var.get(), "function": "PPTMASKMainFunction", "msg": str(e.__class__.__name__), "description": str(e) + "Line No:" + str(e.__traceback__.tb_lineno)})
            raise Exception(e)